"""Assemble deck/Cell_Design_Engine.pptx: each slide PNG full-bleed, with the two
clips embedded over their reserved boxes on slides 5 and 6."""
import json

import imageio
from pptx import Presentation
from pptx.util import Emu

EMU = 9525  # per CSS px (96 dpi); 1280px -> 13.333in
SW, SH = Emu(12192000), Emu(6858000)  # 13.333 x 7.5 in

boxes = json.load(open("deck/video_boxes.json"))


def poster(mp4, out):
    last = None
    rd = imageio.get_reader(mp4)
    for fr in rd:
        last = fr
    rd.close()
    imageio.imwrite(out, last)
    return out


ps = poster("deck/assets/search.mp4", "deck/assets/search_poster.png")
pt = poster("deck/assets/term.mp4", "deck/assets/term_poster.png")

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
blank = prs.slide_layouts[6]

movies = {5: ("search", "deck/assets/search.mp4", ps), 6: ("term", "deck/assets/term.mp4", pt)}

for n in range(1, 11):
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(f"deck/slides/{n:02d}.png", 0, 0, width=SW, height=SH)
    if n in movies:
        key, mp4, pos = movies[n]
        bx = boxes[key]
        slide.shapes.add_movie(mp4, Emu(round(bx["x"] * EMU)), Emu(round(bx["y"] * EMU)),
                               Emu(round(bx["w"] * EMU)), Emu(round(bx["h"] * EMU)),
                               poster_frame_image=pos, mime_type="video/mp4")

prs.save("deck/Cell_Design_Engine.pptx")
print("wrote deck/Cell_Design_Engine.pptx")
